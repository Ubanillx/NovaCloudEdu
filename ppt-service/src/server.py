"""
PPT Template Server：提供 MCP (JSON-RPC) + HTTP REST 双协议接口。
供 Java 后端调用，支持模板解析、图文填充、PPTX 生成。

模板存储由 Java + OSS 管理，Python 侧不再存储模板。

HTTP REST 接口：
  POST /api/templates/parse        解析模板（从 URL 下载）+ 截封面上传 OSS
  POST /api/generate               基于模板生成 PPTX 并上传 OSS

MCP 工具（/mcp 端点）：
  - parse_template
  - generate_ppt_from_template
"""
from __future__ import annotations

import json
import logging

from fastapi import FastAPI, Request
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import Optional

from .schemas import TemplateGenerationSpec, FilledSlide, SlotFill
from .template_registry import registry
from .template_renderer import TemplateRenderer, fill_slide_standalone
from .oss_client import upload_bytes as oss_upload
from .slide_renderer import render_pptx_cover, render_pptx_to_images, render_single_slide, render_pptx_first_page
from .thumbnail import generate_cover as _pillow_cover, render_all_slides as _pillow_render
from .content_validator import validate_slide_fills, format_validation_feedback
from .template_vision_analyzer import build_template_vision_profile, format_vision_profile_for_agent
from .vision_enricher import enrich_template

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
)
logger = logging.getLogger("ppt-mcp-server")

app = FastAPI(
    title="PPT Template MCP Server",
    version="3.0.0",
)


# ==================== MCP 工具定义 ====================

TOOLS = [
    {
        "name": "parse_template",
        "description": (
            "解析用户上传的 PPTX 模板。"
            "传入模板文件下载 URL，"
            "下载并解析后返回完整模板配置 JSON，"
            "包含每页幻灯片的文本槽位信息。"
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "template_url": {
                    "type": "string",
                    "description": "模板 PPTX 文件下载 URL",
                },
            },
            "required": ["template_url"],
        },
    },
    {
        "name": "generate_ppt_from_template",
        "description": (
            "基于模板生成 PPT。"
            "传入模板 URL、标题、"
            "以及每页的克隆来源和填充内容。\n"
            "生成的 PPTX 自动上传到 OSS，"
            "返回文件下载 URL 和前端预览 JSON。"
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "template_url": {
                    "type": "string",
                    "description": "模板 PPTX 文件 URL",
                },
                "title": {
                    "type": "string",
                    "description": "PPT 标题",
                },
                "author": {
                    "type": "string",
                    "description": "作者名称",
                },
                "slides": {
                    "type": "array",
                    "description": "每页幻灯片的克隆来源和填充",
                    "items": {
                        "type": "object",
                        "properties": {
                            "template_slide_index": {
                                "type": "integer",
                                "description": (
                                    "要克隆的模板幻灯片索引"
                                ),
                            },
                            "fills": {
                                "type": "array",
                                "description": (
                                    "槽位填充列表"
                                ),
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "shape_id": {
                                            "type": "integer",
                                            "description": (
                                                "目标形状 ID"
                                            ),
                                        },
                                        "text": {
                                            "type": "string",
                                            "description": (
                                                "填充文本"
                                            ),
                                        },
                                        "items": {
                                            "type": "array",
                                            "items": {
                                                "type": "string",
                                            },
                                            "description": (
                                                "列表项"
                                            ),
                                        },
                                        "image_url": {
                                            "type": "string",
                                            "description": (
                                                "图片 URL（替换图片形状）"
                                            ),
                                        },
                                    },
                                    "required": [
                                        "shape_id",
                                    ],
                                },
                            },
                            "notes": {
                                "type": "string",
                                "description": (
                                    "演讲者备注"
                                ),
                            },
                        },
                        "required": [
                            "template_slide_index",
                            "fills",
                        ],
                    },
                },
            },
            "required": [
                "template_url",
                "title",
                "slides",
            ],
        },
    },
]

SERVER_INFO = {
    "name": "ppt-template-server",
    "version": "3.0.0",
}


# ==================== 核心业务逻辑（MCP + HTTP 共用）====================

def _do_parse_template(template_url: str) -> dict:
    """从 URL 下载并解析模板，截取封面上传 OSS"""
    if not template_url:
        raise ValueError("缺少 template_url 参数")
    cfg = registry.register_from_url(template_url)
    result = cfg.model_dump()

    # 截取封面缩略图并上传 OSS（优先 Gotenberg 高保真，失败退化 Pillow）
    try:
        tpath = registry.get_path(cfg.template_id)
        if tpath:
            try:
                cover_png = render_pptx_cover(tpath)
            except Exception as ge:
                logger.warning("Gotenberg 封面渲染失败，退化 Pillow: %s", ge)
                cover_png = _pillow_cover(tpath)
            cover_url = oss_upload(
                cover_png, ".png", "ppt/cover")
            result["cover_url"] = cover_url
    except Exception as e:
        logger.warning("封面截图失败: %s", e)
        result["cover_url"] = ""

    return result


def _do_generate(arguments: dict) -> dict:
    """基于模板克隆 + 填充生成 PPTX，上传 OSS 返回 URL"""
    turl = arguments.get("template_url")
    if not turl:
        raise ValueError("缺少 template_url 参数")

    cfg = registry.register_from_url(turl)
    tid = cfg.template_id
    tpath = registry.get_path(tid)
    if not cfg or not tpath:
        raise ValueError(f"模板未找到: {tid}")

    spec = TemplateGenerationSpec(**arguments)
    renderer = TemplateRenderer(tpath, spec, cfg)
    pptx_bytes, preview = renderer.render()

    # 上传到 OSS
    file_url = oss_upload(
        pptx_bytes, ".pptx", "ppt/generated")
    title = arguments.get("title", "presentation")

    return {
        "file_name": f"{title}.pptx",
        "file_url": file_url,
        "slide_count": len(arguments.get("slides", [])),
        "preview": [pv.model_dump() for pv in preview],
    }


# ==================== MCP 包装层 ====================

def _mcp_ok(data: dict) -> dict:
    """构造成功的 MCP tool result"""
    data["success"] = True
    return {
        "content": [{
            "type": "text",
            "text": json.dumps(data, ensure_ascii=False),
        }],
    }


def _mcp_err(message: str) -> dict:
    """构造失败的 MCP tool result"""
    return {
        "content": [{
            "type": "text",
            "text": json.dumps({
                "success": False,
                "message": message,
            }, ensure_ascii=False),
        }],
        "isError": True,
    }


def _call_parse_template(arguments: dict) -> dict:
    try:
        return _mcp_ok(
            _do_parse_template(arguments.get("template_url", "")))
    except Exception as e:
        logger.error("parse_template 失败: %s", e, exc_info=True)
        return _mcp_err(str(e))


def _call_generate_from_template(arguments: dict) -> dict:
    try:
        return _mcp_ok(_do_generate(arguments))
    except Exception as e:
        logger.error(
            "generate_from_template 失败: %s", e, exc_info=True)
        return _mcp_err(str(e))


# ==================== MCP JSON-RPC 处理 ====================

_TOOL_HANDLERS = {
    "parse_template": _call_parse_template,
    "generate_ppt_from_template": _call_generate_from_template,
}


def handle_initialize(params: dict) -> dict:
    return {
        "protocolVersion": "2024-11-05",
        "capabilities": {
            "tools": {"listChanged": False},
        },
        "serverInfo": SERVER_INFO,
    }


def handle_tools_list(params: dict) -> dict:
    return {"tools": TOOLS}


def handle_tools_call(params: dict) -> dict:
    tool_name = params.get("name", "")
    arguments = params.get("arguments", {})
    logger.info("工具调用: name=%s", tool_name)

    handler = _TOOL_HANDLERS.get(tool_name)
    if handler is None:
        return _mcp_err(f"未知工具: {tool_name}")

    return handler(arguments)


MCP_METHODS = {
    "initialize": handle_initialize,
    "tools/list": handle_tools_list,
    "tools/call": handle_tools_call,
}


# ==================== HTTP REST 接口 ====================


class ParseTemplateRequest(BaseModel):
    """POST /api/templates/parse 请求体"""
    template_url: str


class SlotFillRequest(BaseModel):
    """fill 单个槽位"""
    shape_id: int
    text: Optional[str] = None
    items: Optional[list[str]] = None
    image_url: Optional[str] = None


class FilledSlideRequest(BaseModel):
    """单页填充"""
    template_slide_index: int
    fills: list[SlotFillRequest] = []
    notes: Optional[str] = None


class GenerateRequest(BaseModel):
    """POST /api/generate 请求体"""
    template_url: str
    title: str
    slides: list[FilledSlideRequest]
    author: Optional[str] = None


@app.post("/api/templates/parse")
async def api_parse_template(req: ParseTemplateRequest):
    """HTTP 接口：解析用户上传模板"""
    try:
        data = _do_parse_template(req.template_url)
        return {"success": True, **data}
    except Exception as e:
        logger.error("API parse_template 失败: %s", e, exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"success": False, "message": str(e)},
        )


class RenderSlidesRequest(BaseModel):
    """POST /api/templates/render-slides 请求体"""
    template_url: str


@app.post("/api/templates/render-slides")
async def api_render_slides(req: RenderSlidesRequest):
    """HTTP 接口：将模板每页渲染为 PNG 图片并上传 OSS"""
    try:
        if not req.template_url:
            raise ValueError("缺少 template_url 参数")
        cfg = registry.register_from_url(req.template_url)
        tpath = registry.get_path(cfg.template_id)
        if not tpath:
            raise ValueError("模板文件未找到")

        try:
            png_list = render_pptx_to_images(tpath)
        except Exception as ge:
            logger.warning("Gotenberg 渲染失败，退化 Pillow: %s", ge)
            png_list = _pillow_render(tpath)
        slide_images = []
        for idx, png_bytes in enumerate(png_list):
            url = oss_upload(png_bytes, ".png", "ppt/slides")
            slide_images.append({"index": idx, "image_url": url})

        return {"success": True, "slide_images": slide_images}
    except Exception as e:
        logger.error("API render-slides 失败: %s", e, exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"success": False, "message": str(e)},
        )


class RenderSingleSlideRequest(BaseModel):
    """POST /api/templates/render-slide 请求体"""
    template_url: str
    slide_index: int


@app.post("/api/templates/render-slide")
async def api_render_single_slide(req: RenderSingleSlideRequest):
    """HTTP 接口：渲染模板指定单页为高清 PNG 并上传 OSS"""
    try:
        if not req.template_url:
            raise ValueError("缺少 template_url 参数")
        cfg = registry.register_from_url(req.template_url)
        tpath = registry.get_path(cfg.template_id)
        if not tpath:
            raise ValueError("模板文件未找到")

        png_bytes = render_single_slide(tpath, req.slide_index)
        url = oss_upload(png_bytes, ".png", "ppt/slides")

        return {"success": True, "image_url": url}
    except Exception as e:
        logger.error(
            "API render-slide 失败: %s", e, exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"success": False, "message": str(e)},
        )


class SlidePreviewRequest(BaseModel):
    """POST /api/generate-slide-preview 请求体"""
    template_url: str
    slide: FilledSlideRequest


def _do_generate_slide_preview(template_url: str, slide_data: dict) -> dict:
    """克隆模板单页 + 填充 + LibreOffice 渲染为 PNG，上传 OSS 返回 URL"""
    import io
    from pptx import Presentation as _Prs
    from .slide_cloner import clone_slide, remove_slide

    if not template_url:
        raise ValueError("缺少 template_url 参数")

    cfg = registry.register_from_url(template_url)
    tpath = registry.get_path(cfg.template_id)
    if not tpath:
        raise ValueError("模板文件未找到")

    # 1. 打开模板
    prs = _Prs(str(tpath))
    orig_count = len(prs.slides)

    # 2. 克隆目标模板页
    tsi = slide_data.get("template_slide_index", 0)
    new_slide = clone_slide(prs, tsi)

    # 3. 填充文本/图片
    filled = FilledSlide(**slide_data)
    fill_slide_standalone(new_slide, filled)

    # 4. 删除所有原始模板页（从后往前）
    for i in range(orig_count - 1, -1, -1):
        remove_slide(prs, i)

    # 5. 保存为 bytes
    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    # 6. 渲染为 PNG（Gotenberg，失败退化 Pillow）
    try:
        png_bytes = render_pptx_first_page(pptx_bytes)
    except Exception as e:
        logger.warning("Gotenberg 单页渲染失败，退化 Pillow: %s", e)
        png_bytes = _pillow_cover(pptx_bytes)

    # 7. 上传 OSS
    image_url = oss_upload(png_bytes, ".png", "ppt/preview")
    return {"image_url": image_url}


@app.post("/api/generate-slide-preview")
async def api_generate_slide_preview(req: SlidePreviewRequest):
    """HTTP 接口：克隆单页 + 填充 + 渲染为 PNG 预览图"""
    try:
        data = _do_generate_slide_preview(
            req.template_url, req.slide.model_dump())
        return {"success": True, **data}
    except Exception as e:
        logger.error("API generate-slide-preview 失败: %s", e, exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"success": False, "message": str(e)},
        )


class ValidateSlideRequest(BaseModel):
    """POST /api/validate-slide 请求体"""
    template_url: str
    slide_index: int = 0
    template_slide_index: int
    fills: list[SlotFillRequest] = []


@app.post("/api/validate-slide")
async def api_validate_slide(req: ValidateSlideRequest):
    """HTTP 接口：校验单页填充内容是否匹配模板 shape 尺寸。

    借鉴 PPTAgent V1 的 _validate_content() 设计思路，
    在 AI 生成内容后、渲染 PPTX 之前进行长度校验，
    避免文字溢出或区域空洞。
    """
    try:
        cfg = registry.register_from_url(req.template_url)

        slide_info = None
        for s in cfg.slides:
            if s.index == req.template_slide_index:
                slide_info = s
                break

        if slide_info is None:
            return JSONResponse(
                status_code=400,
                content={
                    "success": False,
                    "message": (
                        f"模板页索引 {req.template_slide_index}"
                        f" 不存在"
                    ),
                },
            )

        fills = [
            SlotFill(
                shape_id=f.shape_id,
                text=f.text,
                items=f.items,
                image_url=f.image_url,
            )
            for f in req.fills
        ]

        result = validate_slide_fills(
            fills, slide_info, req.slide_index)
        feedback = format_validation_feedback(result)

        return {
            "success": True,
            "is_valid": result.is_valid,
            "issues": [
                i.model_dump() for i in result.issues
            ],
            "suggestions": result.suggestions,
            "feedback_text": feedback,
        }
    except Exception as e:
        logger.error(
            "API validate-slide 失败: %s",
            e, exc_info=True,
        )
        return JSONResponse(
            status_code=500,
            content={
                "success": False,
                "message": str(e),
            },
        )


class AnalyzeTemplateRequest(BaseModel):
    """POST /api/analyze-template 请求体"""
    template_url: str


@app.post("/api/analyze-template")
async def api_analyze_template(req: AnalyzeTemplateRequest):
    """HTTP 接口：Vision-based 模板智能分析。
    借鉴 PPTAgent V1 SlideInducter 设计，对模板每页进行
    语义分析，输出版式分类、适合内容类型、空间分布等。
    """
    try:
        cfg = registry.register_from_url(req.template_url)
        profile = build_template_vision_profile(
            cfg.template_id, cfg.slides,
            cfg.slide_width, cfg.slide_height,
        )
        agent_text = format_vision_profile_for_agent(profile)
        return {
            "success": True,
            "profile": profile.to_dict(),
            "agent_description": agent_text,
        }
    except Exception as e:
        logger.error(
            "API analyze-template 失败: %s",
            e, exc_info=True,
        )
        return JSONResponse(
            status_code=500,
            content={
                "success": False,
                "message": str(e),
            },
        )


class ParseEnrichedRequest(BaseModel):
    """POST /api/templates/parse-enriched 请求体"""
    template_url: str


@app.post("/api/templates/parse-enriched")
async def api_parse_enriched(req: ParseEnrichedRequest):
    """HTTP 接口：解析模板 + 渲染每页 + 多模态语义增强。

    流程：
    1. 下载并解析 PPTX 模板（复用 parse_template）
    2. 渲染每页为 PNG 图片并上传 OSS
    3. 对每页截图 + 结构化 JSON 调用视觉模型分析
    4. 返回 EnrichedTemplateConfig（语义字段 + data 原始数据）
    """
    try:
        if not req.template_url:
            raise ValueError("缺少 template_url 参数")

        # 1. 解析模板
        cfg = registry.register_from_url(req.template_url)
        tpath = registry.get_path(cfg.template_id)
        if not tpath:
            raise ValueError("模板文件未找到")

        # 2. 渲染每页为 PNG
        try:
            png_list = render_pptx_to_images(tpath)
        except Exception as ge:
            logger.warning(
                "Gotenberg 渲染失败，退化 Pillow: %s", ge)
            png_list = _pillow_render(tpath)

        # 3. 上传 PNG 到 OSS
        slide_image_urls: list[str] = []
        for idx, png_bytes in enumerate(png_list):
            url = oss_upload(png_bytes, ".png", "ppt/slides")
            slide_image_urls.append(url)

        # 4. 封面 URL（取第一页）
        cover_url = (
            slide_image_urls[0] if slide_image_urls else ""
        )

        # 5. 多模态语义增强（并发调用视觉模型）
        enriched = await enrich_template(
            cfg, slide_image_urls)
        enriched.cover_url = cover_url

        result = enriched.model_dump()
        result["success"] = True
        return result

    except Exception as e:
        logger.error(
            "API parse-enriched 失败: %s",
            e, exc_info=True,
        )
        return JSONResponse(
            status_code=500,
            content={
                "success": False,
                "message": str(e),
            },
        )


@app.post("/api/generate")
async def api_generate(req: GenerateRequest):
    """HTTP 接口：基于模板生成 PPTX"""
    try:
        data = _do_generate(req.model_dump())
        return {"success": True, **data}
    except Exception as e:
        logger.error("API generate 失败: %s", e, exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"success": False, "message": str(e)},
        )


# ==================== HTML 模式端点 ====================


class HtmlSlideItem(BaseModel):
    """HTML 幻灯片项"""
    slide_html: str
    slide_type: Optional[str] = "content"
    speaker_notes: Optional[str] = ""
    image_suggestions: Optional[list[str]] = []
    generated_image_url: Optional[str] = None
    mode: Optional[str] = "html"


class HtmlGenerateRequest(BaseModel):
    """POST /api/generate-html 请求体"""
    title: str
    slides: list[HtmlSlideItem]
    author: Optional[str] = None
    mode: Optional[str] = "html"


class HtmlSlidePreviewRequest(BaseModel):
    """POST /api/generate-html-slide-preview 请求体"""
    slide_html: str
    generated_image_url: Optional[str] = None


@app.post("/api/generate-html")
async def api_generate_html(req: HtmlGenerateRequest):
    """HTTP 接口：HTML 幻灯片列表 → PPTX 文件"""
    try:
        from .html_renderer import html_slides_to_pptx

        slides_data = [s.model_dump() for s in req.slides]
        pptx_bytes = await html_slides_to_pptx(slides_data, title=req.title)

        file_url = oss_upload(pptx_bytes, ".pptx", "ppt/generated")
        title = req.title or "presentation"

        return {
            "success": True,
            "file_name": f"{title}.pptx",
            "file_url": file_url,
            "slide_count": len(req.slides),
        }
    except Exception as e:
        logger.error(
            "API generate-html 失败: %s", e, exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"success": False, "message": str(e)},
        )


@app.post("/api/generate-html-slide-preview")
async def api_generate_html_slide_preview(
    req: HtmlSlidePreviewRequest,
):
    """HTTP 接口：单页 HTML → PNG 预览图"""
    try:
        from .html_renderer import render_html_to_png

        png_bytes = await render_html_to_png(
            req.slide_html,
            image_url=req.generated_image_url,
        )
        image_url = oss_upload(png_bytes, ".png", "ppt/preview")

        return {"success": True, "image_url": image_url}
    except Exception as e:
        logger.error(
            "API generate-html-slide-preview 失败: %s",
            e, exc_info=True,
        )
        return JSONResponse(
            status_code=500,
            content={"success": False, "message": str(e)},
        )


# ==================== MCP JSON-RPC 端点 ====================

@app.post("/mcp")
async def mcp_endpoint(request: Request):
    """MCP JSON-RPC 2.0 单端点"""
    try:
        body = await request.json()
    except Exception:
        return JSONResponse(content={
            "jsonrpc": "2.0",
            "error": {"code": -32700, "message": "Parse error"},
            "id": None,
        })

    req_id = body.get("id")
    method = body.get("method", "")
    params = body.get("params", {})

    logger.info("MCP 请求: method=%s, id=%s", method, req_id)

    # 客户端初始化通知
    if method == "notifications/initialized":
        return JSONResponse(
            content={"jsonrpc": "2.0", "id": req_id})

    handler = MCP_METHODS.get(method)
    if handler is None:
        return JSONResponse(content={
            "jsonrpc": "2.0",
            "error": {
                "code": -32601,
                "message": f"Method not found: {method}",
            },
            "id": req_id,
        })

    try:
        result = handler(params)
        return JSONResponse(content={
            "jsonrpc": "2.0",
            "result": result,
            "id": req_id,
        })
    except Exception as e:
        logger.error("MCP 处理异常: %s", e, exc_info=True)
        return JSONResponse(content={
            "jsonrpc": "2.0",
            "error": {"code": -32603, "message": str(e)},
            "id": req_id,
        })


@app.get("/health")
async def health_check():
    return {"status": "ok", "service": "ppt-template-server"}


if __name__ == "__main__":
    import os
    import uvicorn

    port = int(os.environ.get("PORT", "8100"))
    uvicorn.run(
        "src.server:app",
        host="0.0.0.0",
        port=port,
        reload=True,
    )
