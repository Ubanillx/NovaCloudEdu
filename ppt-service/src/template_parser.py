"""
模板解析器 v2：读取 PPTX 模板的 **实际幻灯片**，
递归提取文本槽位，识别可填充区域和页面角色，
输出 TemplateConfig JSON 供 AI 匹配使用。
"""
from __future__ import annotations

import logging
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .schemas import (
    TemplateConfig, SlideInfo, TextSlot, ImageSlot,
)

logger = logging.getLogger(__name__)

# 用于识别"可填充"占位文本的关键词
_FILLABLE_KW = re.compile(
    r'单击此处|添加标题|添加内容|请输入|输入文字'
)

# 装饰性文本过滤（全英文且过长 → 大概率是装饰）
_DECO_EN = re.compile(r'^[A-Z\s,.\'-]+$')

# 章节标识
_SECTION_KW = re.compile(r'^PART\s*\d+$', re.I)


# ------------------------------------------------------------------
# 公共 API
# ------------------------------------------------------------------

def parse_template(
    pptx_path: str | Path,
    template_id: str | None = None,
    template_name: str | None = None,
) -> TemplateConfig:
    """解析 PPTX 模板文件，返回 TemplateConfig。"""
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"模板不存在: {path}")

    prs = Presentation(str(path))

    if template_id is None:
        template_id = path.stem
    if template_name is None:
        template_name = path.stem

    slides: list[SlideInfo] = []

    for idx, slide in enumerate(prs.slides):
        text_slots = _extract_text_slots(slide)
        image_slots = _extract_image_slots(slide)
        role = _classify_role(idx, slide, text_slots,
                              len(prs.slides))
        fillable = sum(1 for s in text_slots
                       if s.is_fillable)
        slides.append(SlideInfo(
            index=idx,
            role=role,
            layout_name=slide.slide_layout.name,
            text_slots=text_slots,
            image_slots=image_slots,
            fillable_count=fillable,
        ))

    config = TemplateConfig(
        template_id=template_id,
        name=template_name,
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        slide_count=len(prs.slides),
        slides=slides,
    )

    logger.info(
        "模板解析完成: id=%s, 页数=%d",
        template_id, len(slides),
    )
    return config


# ------------------------------------------------------------------
# 内部：递归提取文本槽位
# ------------------------------------------------------------------

def _extract_text_slots(slide) -> list[TextSlot]:
    """从幻灯片的所有形状中递归提取文本区域。"""
    slots: list[TextSlot] = []
    for shape in slide.shapes:
        _collect_text_shapes(shape, slots)
    return slots


def _extract_image_slots(slide) -> list[ImageSlot]:
    """从幻灯片的所有形状中递归提取图片区域。"""
    slots: list[ImageSlot] = []
    for shape in slide.shapes:
        _collect_image_shapes(shape, slots)
    return slots


def _collect_image_shapes(shape, out: list[ImageSlot]):
    """递归收集图片形状。"""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _collect_image_shapes(child, out)
        return

    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return

    # 跳过面积极小的装饰图片（< 1cm x 1cm）
    MIN_SIZE = 360000  # 1cm in EMU
    w = shape.width or 0
    h = shape.height or 0
    if w < MIN_SIZE or h < MIN_SIZE:
        return

    desc = ""
    if hasattr(shape, "image") and shape.image:
        desc = getattr(shape, "name", "") or ""

    out.append(ImageSlot(
        shape_id=shape.shape_id,
        name=shape.name,
        desc=desc,
        left=shape.left or 0,
        top=shape.top or 0,
        width=w,
        height=h,
    ))


def _collect_text_shapes(shape, out: list[TextSlot]):
    """递归收集含文本的形状。"""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _collect_text_shapes(child, out)
        return

    if not hasattr(shape, 'text_frame'):
        return
    text = shape.text_frame.text.strip()
    if not text:
        return

    # 跳过纯装饰性英文长文本
    if _DECO_EN.match(text) and len(text) > 30:
        return

    is_fillable = bool(_FILLABLE_KW.search(text))
    role = _guess_slot_role(shape, text, is_fillable)

    out.append(TextSlot(
        shape_id=shape.shape_id,
        name=shape.name,
        sample_text=text[:120],
        is_fillable=is_fillable,
        role=role,
        left=shape.left or 0,
        top=shape.top or 0,
        width=shape.width or 0,
        height=shape.height or 0,
    ))


def _guess_slot_role(shape, text: str,
                     is_fillable: bool) -> str:
    """猜测文本槽位的语义角色。"""
    if is_fillable:
        return "body"
    if _SECTION_KW.match(text):
        return "section_number"
    # 短文本 + 不可填充 → 标签/小标题
    if len(text) <= 8:
        return "label"
    return "title"


# ------------------------------------------------------------------
# 内部：分类页面角色
# ------------------------------------------------------------------

def _classify_role(idx: int, slide, slots, total):
    """根据页面位置和内容推断角色。"""
    if idx == 0:
        return "cover"

    # 检查是否章节页
    has_part = any(
        s.role == "section_number" for s in slots
    )
    if has_part:
        return "section"

    # 最后几页
    if idx >= total - 3:
        text_all = ' '.join(s.sample_text for s in slots)
        if '谢谢' in text_all or '感谢' in text_all:
            return "ending"
        if '素材' in text_all or '授权' in text_all:
            return "credits"
        if 'ppt' in text_all.lower():
            return "credits"

    # 有可填充区域 → 内容页
    fillable = sum(1 for s in slots if s.is_fillable)
    if fillable > 0:
        return "content"

    # 第 1 页（非封面非章节）可能是目录
    if idx == 1:
        return "toc"

    return "content"
