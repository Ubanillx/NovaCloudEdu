"""
模板渲染器 v2：克隆模板幻灯片 → 填充文本/图片 → 删除原始页，
保留模板的完整视觉设计（背景、装饰图片、配色等）。
"""
from __future__ import annotations

import io
import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .schemas import (
    TemplateGenerationSpec, FilledSlide,
    PreviewSlide, PreviewElement,
    TemplateConfig,
)
from .slide_cloner import clone_slide, remove_slide
from .utils import download_image_sync
from .template_parser import _FILLABLE_KW
from PIL import Image as _PILImage

logger = logging.getLogger(__name__)


class TemplateRenderer:
    """基于克隆的 PPT 渲染器——保留模板完整视觉"""

    def __init__(
        self,
        template_path: Path,
        spec: TemplateGenerationSpec,
        config: TemplateConfig,
    ):
        self.template_path = template_path
        self.spec = spec
        self.config = config
        self.prs = Presentation(str(template_path))
        self.orig_count = len(self.prs.slides)

    # ----------------------------------------------------------
    # 公共 API
    # ----------------------------------------------------------

    def render(self) -> tuple[bytes, list[PreviewSlide]]:
        """渲染所有幻灯片，返回 (pptx_bytes, preview)。"""
        sw = self.prs.slide_width
        sh = self.prs.slide_height
        preview: list[PreviewSlide] = []

        # 1. 为每条 spec 克隆模板页并追加到末尾
        for out_idx, filled in enumerate(self.spec.slides):
            new_slide = clone_slide(
                self.prs, filled.template_slide_index)
            # 2. 填充文本
            self._fill_slide(new_slide, filled)
            # 3. 备注
            if filled.notes:
                ns = new_slide.notes_slide
                ns.notes_text_frame.text = filled.notes
            # 4. 构建预览
            pv = self._build_preview(
                out_idx, filled, new_slide, sw, sh)
            preview.append(pv)

        # 5. 删除所有原始模板页（从后往前）
        for i in range(self.orig_count - 1, -1, -1):
            remove_slide(self.prs, i)

        buf = io.BytesIO()
        self.prs.save(buf)
        return buf.getvalue(), preview

    # ----------------------------------------------------------
    # 文本填充
    # ----------------------------------------------------------

    def _fill_slide(self, slide, filled: FilledSlide):
        """在克隆后的幻灯片中按 shape_id 填充文本/图片。"""
        # 构建 shape_id → fill 映射
        fill_map = {f.shape_id: f for f in filled.fills}
        # 递归遍历所有形状进行填充
        for shape in slide.shapes:
            self._fill_shape(shape, fill_map, slide)
        # 清除未被填充的占位符文本（防止模板原始文字残留导致重叠）
        for shape in slide.shapes:
            _clear_unfilled_placeholder(shape, fill_map)

    def _fill_shape(self, shape, fill_map: dict, slide):
        """递归查找并填充目标形状（文本或图片）。"""
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                self._fill_shape(child, fill_map, slide)
            return

        fill = fill_map.get(shape.shape_id)
        if fill is None:
            return

        # 图片填充
        if fill.image_url:
            self._replace_image(shape, fill.image_url, slide)
            return

        # 文本填充
        if not hasattr(shape, 'text_frame'):
            return

        tf = shape.text_frame

        if fill.items:
            self._write_bullets(tf, fill.items)
        elif fill.text is not None:
            self._write_text(tf, fill.text)

    def _replace_image(self, shape, image_url: str, slide):
        """下载图片并替换指定形状的图片内容，保持宽高比居中放置。"""
        img_bytes = download_image_sync(image_url)
        if not img_bytes:
            logger.warning(
                "Image download failed, skipping: shape_id=%d url=%s",
                shape.shape_id, image_url)
            return

        # 如果是 Picture 形状，直接替换图片数据（保留模板定义的尺寸）
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                shape.image.blob = img_bytes
                return
            except Exception:
                pass

        # 否则用 add_picture 在同位置插入图片（保持宽高比，居中放置）
        try:
            left, top, w, h = _calc_contain_rect(
                img_bytes, shape.left, shape.top,
                shape.width, shape.height)
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(
                img_stream, left, top, w, h,
            )
        except Exception as e:
            logger.warning(
                "Image insertion failed: shape_id=%d -> %s",
                shape.shape_id, e)

    @staticmethod
    def _write_text(tf, text: str):
        """替换文本框内容，尽量保留原格式。"""
        if not tf.paragraphs:
            return
        # 保留第一段的格式模板
        first_p = tf.paragraphs[0]
        fmt_run = None
        if first_p.runs:
            fmt_run = first_p.runs[0]

        # 清空所有段落
        tf.clear()
        p = tf.paragraphs[0]

        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i > 0:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            # 复制原始字体格式
            if fmt_run is not None:
                _copy_run_fmt(fmt_run, run)

    @staticmethod
    def _write_bullets(tf, items: list[str]):
        """写入列表项，保留格式。"""
        fmt_run = None
        if tf.paragraphs and tf.paragraphs[0].runs:
            fmt_run = tf.paragraphs[0].runs[0]

        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 \
                else tf.add_paragraph()
            p.level = 0
            run = p.add_run()
            run.text = item
            if fmt_run is not None:
                _copy_run_fmt(fmt_run, run)

    # ----------------------------------------------------------
    # 预览 JSON
    # ----------------------------------------------------------

    def _build_preview(
        self, out_idx, filled, slide, sw, sh,
    ) -> PreviewSlide:
        """基于模板 config 中的槽位信息构建预览。"""
        elements: list[PreviewElement] = []

        # 从 config 获取该模板页的 slot 信息
        si = filled.template_slide_index
        slide_info = None
        for s in self.config.slides:
            if s.index == si:
                slide_info = s
                break

        if slide_info:
            slot_map = {
                ts.shape_id: ts
                for ts in slide_info.text_slots
            }
            for fill in filled.fills:
                ts = slot_map.get(fill.shape_id)
                if not ts:
                    continue
                elements.append(PreviewElement(
                    role=ts.role,
                    text=fill.text,
                    items=fill.items,
                    left_pct=round(
                        ts.left / sw * 100, 2),
                    top_pct=round(
                        ts.top / sh * 100, 2),
                    width_pct=round(
                        ts.width / sw * 100, 2),
                    height_pct=round(
                        ts.height / sh * 100, 2),
                ))

        return PreviewSlide(
            slide_index=out_idx,
            template_slide_index=si,
            elements=elements,
        )


# ------------------------------------------------------------------
# 图片尺寸计算工具
# ------------------------------------------------------------------

def _calc_contain_rect(
    img_bytes: bytes,
    box_left: int, box_top: int,
    box_width: int, box_height: int,
) -> tuple[int, int, int, int]:
    """Calculate contain-fit rect preserving aspect ratio.

    Returns (left, top, width, height) in EMU so that the image
    fits entirely within the bounding box and is centered.
    If the image dimensions cannot be read, falls back to the
    original box dimensions (stretch behaviour).
    """
    try:
        with _PILImage.open(io.BytesIO(img_bytes)) as img:
            img_w, img_h = img.size
    except Exception:
        return box_left, box_top, box_width, box_height

    if img_w <= 0 or img_h <= 0:
        return box_left, box_top, box_width, box_height

    img_ratio = img_w / img_h
    box_ratio = box_width / box_height if box_height else 1

    if img_ratio > box_ratio:
        # image is wider than box -> fit by width
        new_w = box_width
        new_h = int(box_width / img_ratio)
    else:
        # image is taller than box -> fit by height
        new_h = box_height
        new_w = int(box_height * img_ratio)

    # center within the bounding box
    left = box_left + (box_width - new_w) // 2
    top = box_top + (box_height - new_h) // 2
    return left, top, new_w, new_h


# ------------------------------------------------------------------
# 独立填充函数（供外部单页预览使用）
# ------------------------------------------------------------------

def fill_slide_standalone(slide, filled: FilledSlide):
    """在克隆后的幻灯片中按 shape_id 填充文本/图片（无需 TemplateRenderer 实例）。"""
    fill_map = {f.shape_id: f for f in filled.fills}
    for shape in slide.shapes:
        _fill_shape_recursive(shape, fill_map, slide)
    # 清除未被填充的占位符文本（防止模板原始文字残留导致重叠）
    for shape in slide.shapes:
        _clear_unfilled_placeholder(shape, fill_map)


def _fill_shape_recursive(shape, fill_map: dict, slide):
    """递归查找并填充目标形状。"""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _fill_shape_recursive(child, fill_map, slide)
        return

    fill = fill_map.get(shape.shape_id)
    if fill is None:
        return

    if fill.image_url:
        img_bytes = download_image_sync(fill.image_url)
        if img_bytes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    shape.image.blob = img_bytes
                    return
                except Exception:
                    pass
            try:
                left, top, w, h = _calc_contain_rect(
                    img_bytes, shape.left, shape.top,
                    shape.width, shape.height)
                img_stream = io.BytesIO(img_bytes)
                slide.shapes.add_picture(
                    img_stream, left, top, w, h,
                )
            except Exception as e:
                logger.warning(
                    "Image insertion failed: shape_id=%d -> %s",
                    shape.shape_id, e)
        return

    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    if fill.items:
        TemplateRenderer._write_bullets(tf, fill.items)
    elif fill.text is not None:
        TemplateRenderer._write_text(tf, fill.text)


# ------------------------------------------------------------------
# 清理未填充的占位符文本（防止模板原始文字残留）
# ------------------------------------------------------------------

def _clear_unfilled_placeholder(shape, fill_map: dict):
    """递归清除未被 AI 填充的占位符文本。

    只清除满足以下条件的形状：
    1. 不在 fill_map 中（AI 没有显式填充）
    2. 文本内容匹配 _FILLABLE_KW 占位关键词（"单击此处"/"添加标题"等）
       或者文本长度 > 15 字符（较长的样例内容也应清除）
    标签/编号等短文本（≤15字符且不含占位关键词）保留不动。
    """
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _clear_unfilled_placeholder(child, fill_map)
        return

    # 已被填充的 shape 跳过
    if shape.shape_id in fill_map:
        return

    if not hasattr(shape, 'text_frame'):
        return

    text = shape.text_frame.text.strip()
    if not text:
        return

    # 匹配占位关键词 → 一定清除
    should_clear = bool(_FILLABLE_KW.search(text))

    # 较长的非标签文本也清除（模板样例内容）
    if not should_clear and len(text) > 15:
        should_clear = True

    if should_clear:
        shape.text_frame.clear()


# ------------------------------------------------------------------
# 格式复制工具
# ------------------------------------------------------------------

def _copy_run_fmt(src_run, dst_run):
    """将源 run 的字体格式复制到目标 run。"""
    try:
        sf = src_run.font
        df = dst_run.font
        if sf.size is not None:
            df.size = sf.size
        if sf.bold is not None:
            df.bold = sf.bold
        if sf.italic is not None:
            df.italic = sf.italic
        if sf.color and sf.color.rgb:
            df.color.rgb = sf.color.rgb
        if sf.name:
            df.name = sf.name
    except Exception:
        pass
