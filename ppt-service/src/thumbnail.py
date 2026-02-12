"""
PPT 封面缩略图生成器。
从 PPTX 文件中提取第一页信息，用 Pillow 合成简易封面预览图。
"""
from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Union

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

# 缩略图输出尺寸
THUMB_WIDTH = 960
THUMB_HEIGHT = 540

# 默认背景色
BG_COLOR = (245, 245, 250)
TEXT_COLOR = (60, 60, 80)
SHAPE_FILL = (220, 225, 235)
IMAGE_PLACEHOLDER = (200, 210, 225)


def generate_cover(source: Union[str, Path, bytes]) -> bytes:
    """
    从 PPTX 生成第一页封面缩略图（PNG 格式）。

    Args:
        source: PPTX 文件路径或字节数据

    Returns:
        PNG 图片字节数据
    """
    if isinstance(source, bytes):
        prs = Presentation(io.BytesIO(source))
    else:
        prs = Presentation(str(source))

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 计算缩放比
    scale_x = THUMB_WIDTH / slide_w
    scale_y = THUMB_HEIGHT / slide_h

    img = Image.new("RGB", (THUMB_WIDTH, THUMB_HEIGHT), BG_COLOR)
    draw = ImageDraw.Draw(img)

    if not prs.slides:
        return _to_png(img)

    slide = prs.slides[0]

    # 绘制背景形状
    _draw_shapes(slide.shapes, draw, scale_x, scale_y)

    buf = _to_png(img)
    logger.info("封面缩略图生成完成: %d bytes", len(buf))
    return buf


def _draw_shapes(shapes, draw: ImageDraw.Draw,
                 sx: float, sy: float):
    """递归绘制形状到画布"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _draw_shapes(shape.shapes, draw, sx, sy)
            continue

        left = int(shape.left * sx)
        top = int(shape.top * sy)
        width = int(shape.width * sx)
        height = int(shape.height * sy)

        if width <= 0 or height <= 0:
            continue

        # 图片形状
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img_bytes = shape.image.blob
                pic = Image.open(io.BytesIO(img_bytes))
                pic = pic.convert("RGB")
                pic = pic.resize((width, height), Image.LANCZOS)
                # 需要把图片粘贴到 draw 所在的 image 上
                draw._image.paste(pic, (left, top))
            except Exception:
                draw.rectangle(
                    [left, top, left + width, top + height],
                    fill=IMAGE_PLACEHOLDER,
                )
            continue

        # 带文本的形状
        if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
            # 轻微填充背景
            draw.rectangle(
                [left, top, left + width, top + height],
                fill=None,
            )
            text = shape.text_frame.text.strip()
            # 截取前 80 字符
            if len(text) > 80:
                text = text[:77] + "..."
            _draw_text(draw, text, left, top, width, height)
            continue

        # 其他形状（装饰色块）
        if shape.shape_type in (
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
        ):
            draw.rectangle(
                [left, top, left + width, top + height],
                fill=SHAPE_FILL,
                outline=(200, 205, 215),
            )


def _draw_text(draw: ImageDraw.Draw, text: str,
               x: int, y: int, w: int, h: int):
    """在矩形区域内绘制文本"""
    # 尝试合适的字号
    font_size = max(10, min(28, h // 3))
    font = None
    for font_path in (
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/System/Library/Fonts/PingFang.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ):
        try:
            font = ImageFont.truetype(font_path, font_size)
            break
        except Exception:
            continue
    if font is None:
        font = ImageFont.load_default()

    # 简单居中
    bbox = draw.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]

    tx = x + max(0, (w - tw) // 2)
    ty = y + max(0, (h - th) // 2)

    draw.text((tx, ty), text, fill=TEXT_COLOR, font=font)


def render_all_slides(source: Union[str, Path, bytes]) -> list[bytes]:
    """
    将 PPTX 的每一页渲染为 PNG 缩略图。

    Args:
        source: PPTX 文件路径或字节数据

    Returns:
        PNG 图片字节列表，按页码顺序
    """
    if isinstance(source, bytes):
        prs = Presentation(io.BytesIO(source))
    else:
        prs = Presentation(str(source))

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    scale_x = THUMB_WIDTH / slide_w
    scale_y = THUMB_HEIGHT / slide_h

    results: list[bytes] = []
    for idx, slide in enumerate(prs.slides):
        img = Image.new("RGB", (THUMB_WIDTH, THUMB_HEIGHT), BG_COLOR)
        draw = ImageDraw.Draw(img)
        _draw_shapes(slide.shapes, draw, scale_x, scale_y)
        results.append(_to_png(img))
        logger.info("渲染第 %d 页缩略图完成", idx)

    logger.info("全部 %d 页缩略图渲染完成", len(results))
    return results


def _to_png(img: Image.Image) -> bytes:
    """将 PIL Image 转为 PNG bytes"""
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()
