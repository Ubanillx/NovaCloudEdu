"""
幻灯片克隆与删除工具。
用于从模板 PPTX 中克隆指定幻灯片并删除多余页。
"""
from __future__ import annotations

from copy import deepcopy

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn


# ------------------------------------------------------------------
# 公共 API
# ------------------------------------------------------------------

def clone_slide(prs: Presentation, slide_index: int):
    """克隆 *slide_index* 处的幻灯片，追加到末尾。

    返回新创建的 ``Slide`` 对象。
    """
    src = prs.slides[slide_index]
    layout = src.slide_layout

    # 1. 通过官方 API 创建空白页（确保正确注册到 Package）
    new_slide = prs.slides.add_slide(layout)

    # 2. 构建 rId 映射: src_rId -> new_rId
    rid_map = _build_rid_map(src, new_slide)

    # 3. 深拷贝源幻灯片的形状树
    src_sp_tree = deepcopy(src.shapes._spTree)

    # 4. 在拷贝的 XML 中重映射 rId
    _remap_rids(src_sp_tree, rid_map)

    # 5. 清空新幻灯片的自动生成形状
    dst_sp_tree = new_slide.shapes._spTree
    for child in list(dst_sp_tree):
        tag = _local_tag(child)
        if tag in _SHAPE_TAGS:
            dst_sp_tree.remove(child)

    # 6. 把源形状追加到新幻灯片
    for child in list(src_sp_tree):
        tag = _local_tag(child)
        if tag in _SHAPE_TAGS:
            dst_sp_tree.append(child)

    # 7. 复制背景（只复制 <p:bg>，不替换整个 cSld）
    src_cSld = src._element.find(qn('p:cSld'))
    dst_cSld = new_slide._element.find(qn('p:cSld'))
    if src_cSld is not None and dst_cSld is not None:
        src_bg_elem = src_cSld.find(qn('p:bg'))
        dst_bg_elem = dst_cSld.find(qn('p:bg'))
        if src_bg_elem is not None:
            new_bg = deepcopy(src_bg_elem)
            _remap_rids(new_bg, rid_map)
            if dst_bg_elem is not None:
                dst_cSld.replace(dst_bg_elem, new_bg)
            else:
                dst_cSld.insert(0, new_bg)

    return new_slide


def remove_slide(prs: Presentation, slide_index: int):
    """删除 *slide_index* 处的幻灯片。"""
    sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
    sld_elem = sldIdLst[slide_index]
    rId = sld_elem.get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sld_elem)


# ------------------------------------------------------------------
# 内部实现
# ------------------------------------------------------------------

_SHAPE_TAGS = frozenset({
    'sp', 'pic', 'grpSp', 'cxnSp', 'graphicFrame',
})


def _local_tag(elem) -> str:
    tag = elem.tag
    return tag.split('}')[-1] if '}' in tag else tag


def _build_rid_map(src_slide, new_slide) -> dict[str, str]:
    """为 *src_slide* → *new_slide* 的关系建立 rId 映射。"""
    rid_map: dict[str, str] = {}

    # layout 关系映射
    src_lr = _find_rid_by_reltype(src_slide.part, RT.SLIDE_LAYOUT)
    new_lr = _find_rid_by_reltype(new_slide.part, RT.SLIDE_LAYOUT)
    if src_lr and new_lr:
        rid_map[src_lr] = new_lr

    # 非 layout 关系逐个复制
    for rId in src_slide.part.rels:
        rel = src_slide.part.rels[rId]
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue
        if rel.is_external:
            new_rId = new_slide.part.rels._add_relationship(
                rel.reltype, rel.target_ref, True,
            )
        else:
            new_rId = new_slide.part.rels._add_relationship(
                rel.reltype, rel.target_part,
            )
        rid_map[rId] = new_rId

    return rid_map


def _find_rid_by_reltype(part, reltype: str):
    for rId in part.rels:
        if part.rels[rId].reltype == reltype:
            return rId
    return None


def _remap_rids(xml_elem, rid_map: dict[str, str]):
    """递归替换 XML 中所有 rId 引用。"""
    for old_rid, new_rid in rid_map.items():
        if old_rid == new_rid:
            continue
        for elem in xml_elem.iter():
            for attr, val in list(elem.attrib.items()):
                if val == old_rid:
                    elem.set(attr, new_rid)
