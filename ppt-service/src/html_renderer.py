"""
HTML 幻灯片渲染器：使用 Playwright 将 HTML 幻灯片渲染为 PNG，
并使用 python-pptx 将 PNG 图片组装为 PPTX。

流程：
  1. HTML → PNG（Playwright headless Chromium 截图）
  2. PNG → PPTX（python-pptx 创建空白幻灯片 + 全页背景图）
"""
from __future__ import annotations

import io
import logging
import re
from typing import Optional

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu

logger = logging.getLogger(__name__)

# 幻灯片尺寸 (16:9)
SLIDE_WIDTH_PX = 1920
SLIDE_HEIGHT_PX = 1080
SLIDE_WIDTH_EMU = Inches(13.333)   # 标准 16:9 宽度
SLIDE_HEIGHT_EMU = Inches(7.5)     # 标准 16:9 高度

# Playwright 异步浏览器实例（懒初始化，全局复用）
_browser = None
_playwright = None


async def _get_browser():
    """懒初始化 Playwright Chromium 浏览器（async）"""
    global _browser, _playwright
    if _browser is None:
        from playwright.async_api import async_playwright
        _playwright = await async_playwright().start()
        _browser = await _playwright.chromium.launch(
            headless=True,
            args=[
                "--no-sandbox",
                "--disable-gpu",
                "--disable-dev-shm-usage",
            ],
        )
        logger.info("Playwright Chromium 浏览器已启动")
    return _browser


async def render_html_to_png(
    slide_html: str,
    width: int = SLIDE_WIDTH_PX,
    height: int = SLIDE_HEIGHT_PX,
    image_url: Optional[str] = None,
) -> bytes:
    """
    将 HTML 幻灯片渲染为 PNG 图片（async）。

    Args:
        slide_html: 自包含的 HTML 幻灯片内容
        width: 视口宽度（默认 1920）
        height: 视口高度（默认 1080）
        image_url: 可选配图 URL

    Returns:
        PNG 图片字节数据
    """
    browser = await _get_browser()

    if image_url:
        slide_html = _inject_image(slide_html, image_url)

    full_html = (
        '<!DOCTYPE html><html><head>'
        '<meta charset="utf-8"><style>'
        '* { margin:0; padding:0; box-sizing:border-box; }'
        f' body {{ width:{width}px; height:{height}px;'
        ' overflow:hidden; }'
        '</style></head><body>'
        f'{slide_html}'
        '</body></html>'
    )

    page = await browser.new_page(
        viewport={"width": width, "height": height},
    )
    try:
        await page.set_content(
            full_html, wait_until="networkidle",
        )
        png_bytes = await page.screenshot(
            type="png", full_page=False,
        )
        logger.info(
            "HTML→PNG 渲染完成: %d bytes", len(png_bytes),
        )
        return png_bytes
    finally:
        await page.close()


def _inject_image(html: str, image_url: str) -> str:
    """将 image-placeholder div 替换为实际图片"""
    # 替换 data-image-suggestion 的 placeholder div
    pattern = r'<div[^>]*class="image-placeholder"[^>]*>.*?</div>'
    replacement = (
        f'<div style="width:100%;height:100%;background-image:url({image_url});'
        f'background-size:cover;background-position:center;"></div>'
    )
    result = re.sub(pattern, replacement, html, flags=re.DOTALL)
    return result


async def html_slides_to_pptx(
    slides: list[dict],
    title: str = "Presentation",
) -> bytes:
    """
    将 HTML 幻灯片列表转换为 PPTX 文件（async）。
    每页 HTML 先渲染为 PNG，然后作为全页背景图插入。
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU

    blank_layout = prs.slide_layouts[6]

    for i, slide_data in enumerate(slides):
        slide_html = slide_data.get("slide_html", "")
        if not slide_html:
            logger.warning(
                "第%d页缺少 slide_html，跳过", i + 1,
            )
            continue

        image_url = slide_data.get("generated_image_url")

        try:
            png_bytes = await render_html_to_png(
                slide_html, image_url=image_url,
            )
        except Exception as e:
            logger.error(
                "第%d页 HTML 渲染失败: %s", i + 1, e,
            )
            png_bytes = _create_placeholder_png(
                f"Slide {i + 1}\n(Render Failed)",
            )

        slide = prs.slides.add_slide(blank_layout)

        img_stream = io.BytesIO(png_bytes)
        slide.shapes.add_picture(
            img_stream,
            Emu(0), Emu(0),
            SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU,
        )

        notes = slide_data.get("speaker_notes", "")
        if notes:
            tf = slide.notes_slide.notes_text_frame
            tf.text = notes

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _create_placeholder_png(text: str) -> bytes:
    """创建纯色占位 PNG 图片"""
    img = Image.new("RGB", (SLIDE_WIDTH_PX, SLIDE_HEIGHT_PX), color=(30, 30, 60))
    # 简单添加文字（不依赖字体文件）
    from PIL import ImageDraw
    draw = ImageDraw.Draw(img)
    draw.text(
        (SLIDE_WIDTH_PX // 2, SLIDE_HEIGHT_PX // 2),
        text,
        fill=(200, 200, 200),
        anchor="mm",
    )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()
